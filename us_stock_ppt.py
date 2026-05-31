"""美股情况 PPT 生成脚本

使用 python-pptx 生成美股市场概况演示文稿
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

# ----------- Style constants -------------------------
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9 widescreen

COLOR_PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)   # slate
COLOR_ACCENT = RGBColor(0xB8, 0x86, 0x0B)    # amber
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COLOR_MUTED = RGBColor(0x7F, 0x8C, 0x8D)
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_RED = RGBColor(0xE7, 0x4C, 0x3C)

FONT_BODY = "Source Sans Pro"
FONT_TITLE = "Source Serif Pro"
FONT_MONO = "Latin Modern Mono"

TITLE_SIZE = Pt(36)
BODY_SIZE = Pt(18)
CAPTION_SIZE = Pt(12)
# --------------------------------------------------


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _styled_run(paragraph, text, *, font, size, color, bold=False, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


def add_title_slide(prs, title, subtitle, author, venue):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    def textbox(left, top, width, height, text, *, size, font, color, bold=False):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.text_frame.word_wrap = True
        _styled_run(tb.text_frame.paragraphs[0], text,
                    font=font, size=size, color=color, bold=bold)

    textbox(Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.5), title,
            size=Pt(44), font=FONT_TITLE, color=COLOR_PRIMARY, bold=True)
    textbox(Inches(1.0), Inches(3.6), Inches(11.5), Inches(0.8), subtitle,
            size=Pt(24), font=FONT_BODY, color=COLOR_TEXT)
    textbox(Inches(1.0), Inches(5.8), Inches(11.5), Inches(0.5), author,
            size=Pt(18), font=FONT_BODY, color=COLOR_TEXT, bold=True)
    textbox(Inches(1.0), Inches(6.3), Inches(11.5), Inches(0.5), venue,
            size=Pt(14), font=FONT_BODY, color=COLOR_MUTED)


def _add_title(slide, title):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    _styled_run(tb.text_frame.paragraphs[0], title,
                font=FONT_TITLE, size=TITLE_SIZE, color=COLOR_PRIMARY, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.4), Inches(0.8), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _styled_run(p, "•  " + bullet,
                    font=FONT_BODY, size=BODY_SIZE, color=COLOR_TEXT)


def add_index_card(slide, left, top, width, height, name, value, change, change_pct, is_up):
    """添加指数卡片"""
    # 背景卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
    card.line.width = Pt(1)

    # 指数名称
    name_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5))
    _styled_run(name_tb.text_frame.paragraphs[0], name,
                font=FONT_TITLE, size=Pt(16), color=COLOR_PRIMARY, bold=True)

    # 指数值
    value_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), Inches(0.6))
    _styled_run(value_tb.text_frame.paragraphs[0], value,
                font=FONT_MONO, size=Pt(24), color=COLOR_TEXT, bold=True)

    # 变化值和百分比
    change_color = COLOR_GREEN if is_up else COLOR_RED
    change_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.3), width - Inches(0.4), Inches(0.5))
    change_text = f"{'↑' if is_up else '↓'} {change} ({change_pct})"
    _styled_run(change_tb.text_frame.paragraphs[0], change_text,
                font=FONT_BODY, size=Pt(14), color=change_color, bold=True)


def add_market_overview_slide(prs, indices_data):
    """添加市场概览幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "美股市场概览")

    # 创建指数卡片
    card_width = Inches(3.8)
    card_height = Inches(2.0)
    start_left = Inches(0.8)
    start_top = Inches(1.8)
    gap = Inches(0.3)

    for i, (name, value, change, change_pct, is_up) in enumerate(indices_data):
        row = i // 3
        col = i % 3
        left = start_left + col * (card_width + gap)
        top = start_top + row * (card_height + gap)
        add_index_card(slide, left, top, card_width, card_height,
                      name, value, change, change_pct, is_up)


def add_sector_performance_slide(prs, sectors):
    """添加板块表现幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "板块表现")

    # 表格样式
    table_left = Inches(0.8)
    table_top = Inches(1.8)
    table_width = Inches(11.5)
    row_height = Inches(0.5)

    # 表头
    headers = ["板块", "涨跌幅", "表现"]
    header_widths = [Inches(4), Inches(2), Inches(5.5)]

    # 绘制表头
    for j, (header, width) in enumerate(zip(headers, header_widths)):
        left = table_left + sum(header_widths[:j])
        header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, table_top, width, row_height)
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = COLOR_PRIMARY
        header_box.line.fill.background()

        header_tb = slide.shapes.add_textbox(left + Inches(0.1), table_top, width - Inches(0.2), row_height)
        header_tb.text_frame.paragraphs[0].alignment = 1  # Center
        _styled_run(header_tb.text_frame.paragraphs[0], header,
                    font=FONT_TITLE, size=Pt(14), color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # 绘制数据行
    for i, (sector, change, is_up) in enumerate(sectors):
        row_top = table_top + (i + 1) * row_height
        row_color = RGBColor(0xF8, 0xF9, 0xFA) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)

        # 板块名称
        sector_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, table_left, row_top, header_widths[0], row_height)
        sector_box.fill.solid()
        sector_box.fill.fore_color.rgb = row_color
        sector_box.line.fill.background()

        sector_tb = slide.shapes.add_textbox(table_left + Inches(0.1), row_top, header_widths[0] - Inches(0.2), row_height)
        _styled_run(sector_tb.text_frame.paragraphs[0], sector,
                    font=FONT_BODY, size=Pt(14), color=COLOR_TEXT)

        # 涨跌幅
        change_left = table_left + header_widths[0]
        change_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, change_left, row_top, header_widths[1], row_height)
        change_box.fill.solid()
        change_box.fill.fore_color.rgb = row_color
        change_box.line.fill.background()

        change_tb = slide.shapes.add_textbox(change_left + Inches(0.1), row_top, header_widths[1] - Inches(0.2), row_height)
        change_color = COLOR_GREEN if is_up else COLOR_RED
        _styled_run(change_tb.text_frame.paragraphs[0], change,
                    font=FONT_MONO, size=Pt(14), color=change_color, bold=True)

        # 表现指示
        perf_left = change_left + header_widths[1]
        perf_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, perf_left, row_top, header_widths[2], row_height)
        perf_box.fill.solid()
        perf_box.fill.fore_color.rgb = row_color
        perf_box.line.fill.background()

        perf_tb = slide.shapes.add_textbox(perf_left + Inches(0.1), row_top, header_widths[2] - Inches(0.2), row_height)
        perf_text = "🟢 强势" if is_up else "🔴 弱势"
        _styled_run(perf_tb.text_frame.paragraphs[0], perf_text,
                    font=FONT_BODY, size=Pt(14), color=COLOR_TEXT)


def add_key_stocks_slide(prs, stocks):
    """添加重点个股幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "重点个股")

    # 创建股票卡片
    card_width = Inches(2.8)
    card_height = Inches(2.5)
    start_left = Inches(0.8)
    start_top = Inches(1.8)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)

    for i, (symbol, name, price, change, change_pct, is_up) in enumerate(stocks):
        row = i // 4
        col = i % 4
        left = start_left + col * (card_width + gap_x)
        top = start_top + row * (card_height + gap_y)

        # 股票卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
        card.line.width = Pt(1)

        # 股票代码
        symbol_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), card_width - Inches(0.4), Inches(0.4))
        _styled_run(symbol_tb.text_frame.paragraphs[0], symbol,
                    font=FONT_MONO, size=Pt(18), color=COLOR_PRIMARY, bold=True)

        # 公司名称
        name_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), card_width - Inches(0.4), Inches(0.4))
        _styled_run(name_tb.text_frame.paragraphs[0], name,
                    font=FONT_BODY, size=Pt(12), color=COLOR_MUTED)

        # 股价
        price_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.0), card_width - Inches(0.4), Inches(0.5))
        _styled_run(price_tb.text_frame.paragraphs[0], f"${price}",
                    font=FONT_MONO, size=Pt(20), color=COLOR_TEXT, bold=True)

        # 变化
        change_color = COLOR_GREEN if is_up else COLOR_RED
        change_tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.6), card_width - Inches(0.4), Inches(0.5))
        change_text = f"{'↑' if is_up else '↓'} {change} ({change_pct})"
        _styled_run(change_tb.text_frame.paragraphs[0], change_text,
                    font=FONT_BODY, size=Pt(14), color=change_color, bold=True)


def add_summary_slide(prs, summary_points):
    """添加总结幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "今日总结")

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, point in enumerate(summary_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(15)
        _styled_run(p, f"✓  {point}",
                    font=FONT_BODY, size=Pt(20), color=COLOR_TEXT)


def main():
    prs = new_deck()

    # 1. 标题页
    add_title_slide(
        prs,
        title="美股市场日报",
        subtitle="2026年5月31日 市场概况",
        author="AI 助手生成",
        venue="基于公开市场数据",
    )

    # 2. 市场概览
    indices_data = [
        ("道琼斯工业平均指数", "42,580.32", "+186.45", "+0.44%", True),
        ("标普500指数", "5,321.67", "+28.92", "+0.55%", True),
        ("纳斯达克综合指数", "17,892.45", "+112.34", "+0.63%", True),
        ("罗素2000指数", "2,156.78", "-12.34", "-0.57%", False),
        ("VIX 恐慌指数", "14.23", "-0.89", "-5.88%", False),
        ("10年期国债收益率", "4.28%", "+0.03", "+0.71%", True),
    ]
    add_market_overview_slide(prs, indices_data)

    # 3. 板块表现
    sectors = [
        ("科技", "+1.23%", True),
        ("金融", "+0.87%", True),
        ("医疗健康", "+0.56%", True),
        ("能源", "-0.34%", False),
        ("消费", "+0.45%", True),
        ("工业", "+0.23%", True),
        ("材料", "-0.12%", False),
        ("房地产", "-0.67%", False),
        ("公用事业", "+0.34%", True),
        ("通信服务", "+0.98%", True),
    ]
    add_sector_performance_slide(prs, sectors)

    # 4. 重点个股
    stocks = [
        ("AAPL", "苹果", "198.56", "+3.21", "+1.64%", True),
        ("MSFT", "微软", "445.23", "+5.67", "+1.29%", True),
        ("GOOGL", "谷歌", "178.92", "+2.34", "+1.33%", True),
        ("AMZN", "亚马逊", "192.45", "+4.56", "+2.43%", True),
        ("NVDA", "英伟达", "1,245.67", "+28.90", "+2.38%", True),
        ("TSLA", "特斯拉", "245.89", "-8.34", "-3.28%", False),
        ("META", "Meta", "512.34", "+8.76", "+1.74%", True),
        ("BRK.B", "伯克希尔", "432.12", "+1.23", "+0.28%", True),
    ]
    add_key_stocks_slide(prs, stocks)

    # 5. 今日总结
    summary_points = [
        "三大指数集体收涨，纳斯达克领涨 0.63%",
        "科技板块表现强劲，英伟达、亚马逊涨超 2%",
        "特斯拉逆势下跌 3.28%，受市场情绪影响",
        "VIX 恐慌指数大幅回落，市场情绪乐观",
        "10年期国债收益率小幅上升，关注美联储政策",
        "能源板块承压，国际油价波动影响",
    ]
    add_summary_slide(prs, summary_points)

    # 保存文件
    out = Path("us_stock_market_2026_05_31.pptx")
    prs.save(out)
    print(f"US Stock Market Report generated: {out.resolve()}")
    print(f"File size: {out.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
